"""
Document acquisition module.

Handles processing of PDFs, Office documents (DOCX, PPTX, XLSX),
and other document formats.
"""

from __future__ import annotations

import hashlib
import mimetypes
import shutil
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, BinaryIO

from inception.config import get_config


@dataclass
class DocumentInfo:
    """Metadata and content info for a document."""
    
    path: Path
    mime_type: str
    size_bytes: int
    content_hash: str
    
    # Metadata
    title: str | None = None
    author: str | None = None
    created: datetime | None = None
    modified: datetime | None = None
    
    # Structure info
    page_count: int | None = None
    word_count: int | None = None
    has_images: bool = False
    has_tables: bool = False
    
    # Processing state
    processed: bool = False
    extracted_text: str | None = None


@dataclass
class PDFPage:
    """Extracted content from a PDF page."""
    
    page_num: int  # 0-indexed
    text: str | None = None
    tables: list[list[list[str]]] = field(default_factory=list)
    images: list[dict[str, Any]] = field(default_factory=list)
    
    # Layout info
    width: float | None = None
    height: float | None = None
    
    # OCR info (for scanned pages)
    is_scanned: bool = False
    ocr_confidence: float | None = None


@dataclass
class PDFContent:
    """Full extracted content from a PDF."""
    
    path: Path
    info: DocumentInfo
    pages: list[PDFPage] = field(default_factory=list)
    
    # Document-level extraction
    toc: list[dict[str, Any]] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)
    
    @property
    def full_text(self) -> str:
        """Get concatenated text from all pages."""
        return "\n\n".join(
            page.text or "" for page in self.pages
        )
    
    @property
    def page_count(self) -> int:
        return len(self.pages)


@dataclass
class SlideContent:
    """Extracted content from a presentation slide."""
    
    slide_num: int  # 1-indexed
    title: str | None = None
    text: str | None = None
    notes: str | None = None
    
    # Visual elements
    shapes: list[dict[str, Any]] = field(default_factory=list)
    images: list[dict[str, Any]] = field(default_factory=list)
    tables: list[list[list[str]]] = field(default_factory=list)


@dataclass
class PresentationContent:
    """Full extracted content from a presentation."""
    
    path: Path
    info: DocumentInfo
    slides: list[SlideContent] = field(default_factory=list)
    
    @property
    def full_text(self) -> str:
        """Get concatenated text from all slides."""
        parts = []
        for slide in self.slides:
            if slide.title:
                parts.append(f"# {slide.title}")
            if slide.text:
                parts.append(slide.text)
            if slide.notes:
                parts.append(f"\n[Speaker Notes: {slide.notes}]")
        return "\n\n".join(parts)
    
    @property
    def slide_count(self) -> int:
        return len(self.slides)


def get_document_info(path: Path) -> DocumentInfo:
    """
    Get basic information about a document file.
    
    Args:
        path: Path to the document
    
    Returns:
        DocumentInfo with metadata
    """
    path = Path(path)
    
    if not path.exists():
        raise FileNotFoundError(f"Document not found: {path}")
    
    # Get MIME type
    mime_type, _ = mimetypes.guess_type(str(path))
    if mime_type is None:
        mime_type = "application/octet-stream"
    
    # Compute content hash
    with open(path, "rb") as f:
        content_hash = hashlib.sha256(f.read()).hexdigest()
    
    return DocumentInfo(
        path=path,
        mime_type=mime_type,
        size_bytes=path.stat().st_size,
        content_hash=content_hash,
    )


def extract_pdf(
    path: Path,
    extract_tables: bool = True,
    extract_images: bool = False,
    use_ocr: bool = True,
) -> PDFContent:
    """
    Extract content from a PDF file.
    
    Args:
        path: Path to PDF file
        extract_tables: Whether to extract tables
        extract_images: Whether to extract embedded images
        use_ocr: Whether to use OCR for scanned pages
    
    Returns:
        PDFContent with extracted text and structure
    """
    import pdfplumber
    
    path = Path(path)
    info = get_document_info(path)
    
    pages: list[PDFPage] = []
    
    with pdfplumber.open(path) as pdf:
        info.page_count = len(pdf.pages)
        
        # Extract metadata
        metadata = pdf.metadata or {}
        if metadata.get("Title"):
            info.title = metadata["Title"]
        if metadata.get("Author"):
            info.author = metadata["Author"]
        
        for i, page in enumerate(pdf.pages):
            pdf_page = PDFPage(
                page_num=i,
                width=page.width,
                height=page.height,
            )
            
            # Extract text
            text = page.extract_text()
            if text:
                pdf_page.text = text
            elif use_ocr:
                # Mark as potentially needing OCR
                pdf_page.is_scanned = True
            
            # Extract tables
            if extract_tables:
                tables = page.extract_tables()
                if tables:
                    pdf_page.tables = tables
                    info.has_tables = True
            
            # Extract images would go here
            if extract_images:
                # pdfplumber can list images but extracting is complex
                pass
            
            pages.append(pdf_page)
    
    return PDFContent(
        path=path,
        info=info,
        pages=pages,
    )


def extract_pptx(path: Path) -> PresentationContent:
    """
    Extract content from a PowerPoint file.
    
    Args:
        path: Path to PPTX file
    
    Returns:
        PresentationContent with extracted slides
    """
    from pptx import Presentation
    from pptx.util import Inches
    
    path = Path(path)
    info = get_document_info(path)
    
    prs = Presentation(str(path))
    
    # Extract core properties if available
    if prs.core_properties:
        info.title = prs.core_properties.title
        info.author = prs.core_properties.author
        info.created = prs.core_properties.created
        info.modified = prs.core_properties.modified
    
    slides: list[SlideContent] = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_content = SlideContent(slide_num=i)
        
        # Extract title
        if slide.shapes.title:
            slide_content.title = slide.shapes.title.text
        
        # Extract text from all shapes
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                # Skip title (already extracted)
                if hasattr(slide.shapes, "title") and shape == slide.shapes.title:
                    continue
                text_parts.append(shape.text)
            
            # Check for tables
            if shape.has_table:
                info.has_tables = True
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                slide_content.tables.append(table_data)
        
        slide_content.text = "\n".join(text_parts) if text_parts else None
        
        # Extract speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text:
                slide_content.notes = notes_frame.text
        
        slides.append(slide_content)
    
    info.page_count = len(slides)
    
    return PresentationContent(
        path=path,
        info=info,
        slides=slides,
    )


def extract_docx(path: Path) -> tuple[DocumentInfo, str]:
    """
    Extract content from a Word document.
    
    Args:
        path: Path to DOCX file
    
    Returns:
        Tuple of (DocumentInfo, extracted_text)
    """
    from docx import Document
    
    path = Path(path)
    info = get_document_info(path)
    
    doc = Document(str(path))
    
    # Extract core properties
    if doc.core_properties:
        info.title = doc.core_properties.title
        info.author = doc.core_properties.author
        info.created = doc.core_properties.created
        info.modified = doc.core_properties.modified
    
    # Extract text
    paragraphs = []
    for para in doc.paragraphs:
        if para.text:
            paragraphs.append(para.text)
    
    # Extract tables
    for table in doc.tables:
        info.has_tables = True
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            paragraphs.append(row_text)
    
    text = "\n\n".join(paragraphs)
    info.extracted_text = text
    info.processed = True
    
    return info, text


def extract_xlsx(path: Path) -> tuple[DocumentInfo, list[dict[str, Any]]]:
    """
    Extract content from an Excel spreadsheet.
    
    Args:
        path: Path to XLSX file
    
    Returns:
        Tuple of (DocumentInfo, list of sheet data)
    """
    from openpyxl import load_workbook
    
    path = Path(path)
    info = get_document_info(path)
    info.has_tables = True
    
    wb = load_workbook(str(path), read_only=True, data_only=True)
    
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        
        # Read all rows
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append(list(row))
        
        sheets.append({
            "name": sheet_name,
            "rows": rows,
        })
    
    wb.close()
    
    return info, sheets


def copy_to_artifacts(
    source_path: Path,
    preserve_name: bool = True,
) -> Path:
    """
    Copy a document to the artifacts directory.
    
    Args:
        source_path: Path to source document
        preserve_name: Whether to preserve original filename
    
    Returns:
        Path to copied file in artifacts directory
    """
    config = get_config()
    
    artifacts_dir = config.artifacts_dir / "documents"
    artifacts_dir.mkdir(parents=True, exist_ok=True)
    
    source_path = Path(source_path)
    
    if preserve_name:
        dest_path = artifacts_dir / source_path.name
        # Handle duplicates
        if dest_path.exists():
            stem = source_path.stem
            suffix = source_path.suffix
            counter = 1
            while dest_path.exists():
                dest_path = artifacts_dir / f"{stem}_{counter}{suffix}"
                counter += 1
    else:
        # Use content hash as filename
        info = get_document_info(source_path)
        suffix = source_path.suffix
        dest_path = artifacts_dir / f"{info.content_hash[:16]}{suffix}"
    
    shutil.copy2(source_path, dest_path)
    return dest_path


def detect_document_type(path: Path) -> str:
    """
    Detect the type of a document.
    
    Args:
        path: Path to document
    
    Returns:
        Document type string: 'pdf', 'pptx', 'docx', 'xlsx', 'unknown'
    """
    path = Path(path)
    
    # First try by extension
    ext = path.suffix.lower()
    ext_map = {
        ".pdf": "pdf",
        ".pptx": "pptx",
        ".ppt": "ppt",
        ".docx": "docx",
        ".doc": "doc",
        ".xlsx": "xlsx",
        ".xls": "xls",
        ".epub": "epub",
        ".txt": "txt",
        ".md": "markdown",
        ".html": "html",
        ".htm": "html",
    }
    
    if ext in ext_map:
        return ext_map[ext]
    
    # Try by MIME type
    mime_type, _ = mimetypes.guess_type(str(path))
    if mime_type:
        if "pdf" in mime_type:
            return "pdf"
        if "powerpoint" in mime_type or "presentation" in mime_type:
            return "pptx"
        if "word" in mime_type or "document" in mime_type:
            return "docx"
        if "excel" in mime_type or "spreadsheet" in mime_type:
            return "xlsx"
    
    return "unknown"
